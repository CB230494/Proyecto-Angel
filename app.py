# =========================
# 📊 MPGP – Proyecto INTEGRADO (todo en un diagrama)
# =========================
# Principal + Nodos Demandantes (Proc. 1.4 parte 1) + Conducta Delictiva Reiterada (3 carriles)
# Conexiones:
#   SÍ-1 → Nodos → SÍ-3
#   SÍ-6 → Reiterada → SÍ-7
# Exporta PNG / PDF / PPTX.
# =========================

import io, math
from typing import List, Tuple
from PIL import Image, ImageDraw, ImageFont, Image as PILImage
import streamlit as st
from pptx import Presentation
from pptx.util import Inches

# ---------- Config / Estilos ----------
st.set_page_config(page_title="MPGP – Proyecto Integrado", layout="wide")
st.title("Modelo Preventivo de Gestión Policial – Función de Operacionales (Proyecto Integrado)")
st.caption("Diagrama principal + Nodos Demandantes + Conducta Delictiva Reiterada en un solo flujo. Exporta PNG / PDF / PPTX.")

BLUE=(31,78,121); BORDER=(155,187,217); LIGHTBLUE=(220,235,247); LIGHTYELLOW=(255,248,225)
LANE_BG=(234,240,249)
WHITE=(255,255,255); BLACK=(0,0,0); BG=(247,250,255)

W=2000
BASE_H=1400
H_NODOS=760
H_REIT =920
SAFE=16
ARROW_HEAD=18
HEAD_CLEAR=SAFE+ARROW_HEAD+6

def font(size:int):
    try: return ImageFont.truetype("DejaVuSans.ttf", size)
    except: return ImageFont.load_default()
FONT=font(22); FONT_SMALL=font(18); FONT_TITLE=font(28); FONT_LANE=font(20); FONT_SUB=font(24)

# ---------- Entradas (todos los textos) ----------
# Principal
cA,cB = st.columns(2)
with cA:
    t_inicio=st.text_input("INICIO","Planificación preventiva anual")
    b1=st.text_area("Bloque 1","Definición y calendarización de Delegaciones\n(Procedimiento 1.1 MPGP)")
    b2=st.text_area("Bloque 2","Apreciación situacional del territorio\n(Procedimiento 1.2)")
    b3=st.text_area("Bloque 3","Identificación de factores de riesgo y delitos\n(DATAPOL, estadísticas, patrullaje)")
    q_dec=st.text_input("Decisión","¿Se identifican riesgos prioritarios?")
with cB:
    s1=st.text_area("SÍ 1","Priorización de riesgos y delitos\n(Pareto, MIC-MAC, Triángulo de violencias)")
    s2=st.text_area("SÍ 2","Construcción de líneas de acción preventivas\n(Procedimiento 2.3)")
    s3=st.text_area("SÍ 3","Planificación de programas policiales preventivos\n(Procedimiento 2.4)")
    s4=st.text_area("SÍ 4","Elaboración de órdenes de servicio para operativos")
    s5=st.text_area("SÍ 5","Implementación en terreno\n• Patrullajes preventivos\n• Respuesta inmediata\n• Supervisión\n• Coordinación local")
    s6=st.text_area("SÍ 6","Reporte de operativos (RAP, DATAPOL, informes)")
    s7=st.text_area("SÍ 7","Evaluación de cumplimiento (Trazabilidad 3.1 y 3.2)")
    s8=st.text_area("SÍ 8","Retroalimentación a la planificación preventiva")
cC,cD = st.columns(2)
with cC:
    n1=st.text_area("NO 1","Patrullaje rutinario y vigilancia continua")
    n2=st.text_area("NO 2","Registro de factores menores en RAP")
    n3=st.text_area("NO 3","Integración al análisis situacional")
with cD:
    t_fin=st.text_area("FIN","Evaluación global de resultados\n(Indicadores, metas, impacto – 3.3)")

st.markdown("### ⚙️ Ajustes del principal")
g1,g2,g3 = st.columns(3)
with g1:
    start_offset=st.slider("Inicio rama SÍ (relativo al rombo)", -800, 240, -280, 5)
with g2:
    step_user=st.slider("Paso objetivo (px)", 90, 220, 130, 5)
with g3:
    altura_si = st.slider("Altura cajas SÍ (px)", 70, 110, 78, 2)
g4,g5,g6 = st.columns(3)
with g4:
    altura_si5 = st.slider("Altura cuadro SÍ 5 (px)", 110, 260, 150, 5)
with g5:
    ancho_si   = st.slider("Ancho cuadros rama SÍ (px)", 520, 680, 560, 10)
with g6:
    retro_rail = st.slider("Separación lateral retroalimentación (px)", 120, 260, 165, 5)
c7,c8,c9 = st.columns(3)
with c7:
    comp_branch = st.slider("Compacidad rama SÍ (0.20–1.00)", 0.20, 1.00, 0.45, 0.05)
with c8:
    comp_early  = st.slider("Compacidad extra SÍ 1–5", 0.20, 1.00, 0.35, 0.05)
with c9:
    vert_margin = st.slider("Margen vertical flecha (px)", 24, 70, 30, 2)

# Subproceso INTEGRADO A: Nodos Demandantes (Proc. 1.4 parte 1)
st.markdown("### 🧩 Subproceso: Nodos Demandantes (Proc. 1.4 – parte 1)")
a_sup1 = st.text_area("A-Superior 1","Convoca a reunión EDO de segundo nivel con plantillas diferenciadas")
a_sup2 = st.text_area("A-Superior 2","Verifica insumos mínimos para análisis (capas, encuestas, informes)")
a_sup3 = st.text_area("A-Superior 3","Completa la Matriz de Nodos demandantes priorizados")
a_sup4 = st.text_area("A-Superior 4","Elabora órdenes de servicio para evidencia y monitoreo")
a_inf1 = st.text_area("A-Inferior 1","Abre el SIG para visualizar y mapear la información de nodos")
a_inf2 = st.text_area("A-Inferior 2","Selecciona capas de información disponibles en el SIG")
a_inf3 = st.text_area("A-Inferior 3","Presenta factores de riesgo críticos y variaciones del mes anterior")
a_inf4 = st.text_area("A-Inferior 4","Analiza puntos críticos y oportunidades (Análisis cualitativo)")
a_inf5 = st.text_area("A-Inferior 5","Diseña respuestas policiales diferenciadas (prevención/disuasión)")
a_inf6 = st.text_area("A-Inferior 6","Presenta avance de cumplimiento de órdenes y coordinación interinstitucional")
anx_wA = st.slider("Ancho de cajas Nodos (px)", 360, 560, 420, 10)
anx_gapA = st.slider("Separación horizontal Nodos (px)", 40, 140, 70, 5)

# Subproceso INTEGRADO B: Conducta Delictiva Reiterada (swimlanes)
st.markdown("### 🧩 Subproceso: Conducta Delictiva Reiterada (3 carriles)")
lane1 = st.text_input("Carril 1 (arriba)","Asesor(a) legal de la Dirección Regional")
lane2 = st.text_input("Carril 2 (medio)","Oficial de Operaciones de la Dirección Regional")
lane3 = st.text_input("Carril 3 (abajo)","Agente de Operaciones de la Delegación Policial")
r1 = st.text_area("B-Superior 1","Realiza el estudio de antecedentes judiciales a cada objetivo priorizado ante el Ministerio Público")
r2 = st.text_area("B-Superior 2","Elabora la ficha de personas con conducta delictiva reiterada")
r3 = st.text_area("B-Superior 3","Remite las fichas a la oficina de operaciones regional para su distribución")
r4 = st.text_area("B-Superior 4","Participa y presenta las fichas en la reunión EDO de Planificación/Testeo (primer nivel)")
rm1 = st.text_area("B-Medio 1","Envía las fichas a las oficinas de operaciones de las Delegaciones Policiales")
rm2 = st.text_area("B-Medio 2","Incluye las fichas como documentación para la reunión EDO (primer nivel)")
rb1 = st.text_area("B-Inferior 1","Incluye las fichas como documentación para la reunión EDO (segundo nivel)")
anx_wB = st.slider("Ancho de cajas Reiterada (px)", 360, 520, 420, 10)
anx_gapB = st.slider("Separación horizontal Reiterada (px)", 40, 120, 70, 5)

# ---------- Utilidades de dibujo (ROBUSTECIDAS) ----------
def _as_text(x) -> str:
    """Devuelve siempre un string ('' si x es None)."""
    return "" if x is None else str(x)

def wrap_text(d: ImageDraw.ImageDraw, text, font: ImageFont.FreeTypeFont, max_w: int) -> List[str]:
    text = _as_text(text)
    out=[]
    for raw in text.split("\n"):
        words=raw.split(" "); line=""
        for w in words:
            t=(line+" "+w).strip()
            if d.textlength(t,font=font)<=max_w: line=t
            else:
                if line: out.append(line); line=w
        if line: out.append(line)
    return out

def draw_centered(d: ImageDraw.ImageDraw, text, box, font=FONT, fill=BLACK, leading=6):
    text = _as_text(text)
    x0,y0,x1,y1=box; max_w=x1-x0-30
    lines=wrap_text(d,text,font,max_w); lh=font.size+leading; total=len(lines)*lh
    y=y0+(y1-y0-total)//2
    for ln in lines:
        ln=_as_text(ln)
        w=d.textlength(ln,font=font); x=x0+(x1-x0-w)//2
        d.text((x,y),ln,font=font,fill=fill); y+=lh

def rrect(d, box, radius=22, fill=WHITE, outline=BLUE, width=3): d.rounded_rectangle(box,radius=radius,fill=fill,outline=outline,width=width)
def oval(d, box, fill=LIGHTBLUE, outline=BLUE, width=3): d.ellipse(box,fill=fill,outline=outline,width=width)
def diamond(d, box, fill=LIGHTYELLOW, outline=BLUE, width=3):
    x0,y0,x1,y1=box; cx=(x0+x1)//2; cy=(y0+y1)//2
    d.polygon([(cx,y0),(x1,cy),(cx,y1),(x0,cy)], fill=fill, outline=outline)

def arrow(d: ImageDraw.ImageDraw, p1: Tuple[int,int], p2: Tuple[int,int], color=BLUE, width=4):
    d.line([p1,p2], fill=color, width=width)
    ang=math.atan2(p2[1]-p1[1], p2[0]-p1[0])
    a1=(p2[0]-ARROW_HEAD*math.cos(ang-0.4), p2[1]-ARROW_HEAD*math.sin(ang-0.4))
    a2=(p2[0]-ARROW_HEAD*math.cos(ang+0.4), p2[1]-ARROW_HEAD*math.sin(ang+0.4))
    d.polygon([p2,a1,a2], fill=color)

def arrow_down(d: ImageDraw.ImageDraw, p1: Tuple[int,int], p2: Tuple[int,int], **kw):
    if p2[1] < p1[1]: p1, p2 = p2, p1
    arrow(d, p1, p2, **kw)

def poly_arrow(d: ImageDraw.ImageDraw, pts, color=BLUE, width=4):
    for i in range(len(pts)-2): d.line([pts[i], pts[i+1]], fill=color, width=width)
    arrow(d, pts[-2], pts[-1], color=color, width=width)

def paste_vertical_label(img: PILImage.Image, box: Tuple[int,int,int,int], text: str, bg=LANE_BG, fg=BLUE):
    text=_as_text(text)
    x0,y0,x1,y1 = box
    draw = ImageDraw.Draw(img)
    draw.rectangle(box, fill=bg, outline=BORDER, width=2)
    tmp = PILImage.new("RGBA", (y1-y0, x1-x0), (0,0,0,0))
    td = ImageDraw.Draw(tmp)
    tw = td.textlength(text, font=FONT_LANE)
    td.text((((y1-y0)-tw)//2, ((x1-x0)-FONT_LANE.size)//2), text, font=FONT_LANE, fill=fg)
    rot = tmp.rotate(90, expand=True)
    img.paste(rot, (x0, y0), rot)

# ---------- Render ----------
def render_png() -> bytes:
    HH = BASE_H + H_NODOS + H_REIT
    img=PILImage.new("RGB",(W,HH),BG); d=ImageDraw.Draw(img)
    d.rectangle([20,20,W-20,HH-20], outline=BORDER, width=3)
    d.text((W//2,50),"Modelo Preventivo de Gestión Policial – Función de Operacionales (Proyecto Integrado)",
           font=FONT_TITLE, fill=BLUE, anchor="mm")

    # ===== Principal =====
    cx=W//2; vgap=130; bw,bh=480,104; y0=120
    def box(x,y,w=bw,h=bh): return [x-w//2, y-h//2, x+w//2, y+h//2]
    def big(x,y,w,h):       return [x-w//2, y-h//2, x+w//2, y+h//2]

    r_inicio=box(cx,y0);        oval(d,r_inicio); draw_centered(d,f"INICIO\n{t_inicio}",r_inicio)
    r1=box(cx,y0+vgap);         rrect(d,r1); draw_centered(d,b1,r1)
    r2=box(cx,y0+vgap*2);       rrect(d,r2); draw_centered(d,b2,r2)
    r3=box(cx,y0+vgap*3);       rrect(d,r3); draw_centered(d,b3,r3)
    r_dec=big(cx,y0+vgap*4,520,124); diamond(d,r_dec); draw_centered(d,q_dec,r_dec)
    r_fin=big(cx,y0+vgap*8+60,560,124); oval(d,r_fin); draw_centered(d,f"FIN\n{t_fin}",r_fin)

    def top_pt(r):   return ((r[0]+r[2])//2, r[1]-HEAD_CLEAR)
    def bot_pt(r):   return ((r[0]+r[2])//2, r[3]+SAFE)
    def left_pt(r):  return (r[0]-HEAD_CLEAR, (r[1]+r[3])//2)
    def right_pt(r): return (r[2]+HEAD_CLEAR, (r[1]+r[3])//2)

    for a,b in [(r_inicio,r1),(r1,r2),(r2,r3),(r3,r_dec)]:
        arrow_down(d, bot_pt(a), top_pt(b))
    arrow_down(d, ((r_dec[0]+r_dec[2])//2, r_dec[3]+SAFE), top_pt(r_fin))

    # ----- Rama SÍ compacta -----
    rx=cx+620; TOP_GUARD=140
    widths=int(ancho_si)
    start_y=max(TOP_GUARD, (r_dec[1]+r_dec[3])//2 + start_offset)
    heights=[int(altura_si), int(altura_si), int(altura_si), int(altura_si),
             int(altura_si5), int(altura_si), int(altura_si), int(altura_si)]
    max_center_y = r_fin[1]-40
    req=[(heights[i]/2 + heights[i+1]/2 + vert_margin) for i in range(7)]
    mult=[comp_branch]*7
    for i in range(0,4): mult[i]=min(comp_branch, comp_early)
    available = max_center_y - heights[-1]/2 - start_y
    min_total = sum(req)
    if available < min_total:
        start_y = max(TOP_GUARD, max_center_y - heights[-1]/2 - min_total)
        available = max_center_y - heights[-1]/2 - start_y

    def total_for(bs: float) -> float:
        return sum(max(req[i], bs*mult[i]) for i in range(7))
    lo,hi=0.0,float(step_user)
    for _ in range(30):
        mid=(lo+hi)/2
        if total_for(mid)<=available: lo=mid
        else: hi=mid
    base_step=lo
    steps=[max(req[i], base_step*mult[i]) for i in range(7)]
    Ys=[start_y]
    for s in steps: Ys.append(Ys[-1]+s)
    textos=[s1,s2,s3,s4,s5,s6,s7,s8]; rs=[]
    for i,(t,y) in enumerate(zip(textos,Ys)):
        r=big(rx,int(y),widths,heights[i]); rrect(d,r); draw_centered(d,t,r); rs.append(r)

    # ----- Rama NO (alineada con SÍ 1-3)
    lx=cx-620; rn=[]
    for i,t in enumerate([n1,n2,n3]):
        r=big(lx,int(Ys[i]),widths,int(altura_si)); rrect(d,r); draw_centered(d,t,r); rn.append(r)

    # Decisión + rótulos
    seg_si=(right_pt(r_dec),(rs[0][0]-HEAD_CLEAR,(rs[0][1]+rs[0][3])//2))
    seg_no=(left_pt(r_dec), (rn[0][2]+HEAD_CLEAR,(rn[0][1]+rn[0][3])//2))
    arrow(d,*seg_si); arrow(d,*seg_no)
    def label_out(p1,p2,text):
        mx,my=(p1[0]+p2[0])/2,(p1[1]+p2[1])/2
        dx,dy=p2[0]-p1[0],p2[1]-p1[1]; L=max(1.0,math.hypot(dx,dy))
        nx,ny=-dy/L,dx/L
        cand1=(mx+nx*36,my+ny*36); cand2=(mx-nx*36,my-ny*36)
        tx,ty = cand1 if abs(cand1[0]-cx)>abs(cand2[0]-cx) else cand2
        w=d.textlength(text,font=FONT_SMALL); h=FONT_SMALL.size; pad=6
        d.rounded_rectangle([tx-w/2-pad,ty-h/2-pad,tx+w/2+pad,ty+h/2+pad], radius=8, fill=WHITE)
        d.text((tx,ty), text, font=FONT_SMALL, fill=BLUE, anchor="mm")
    label_out(*seg_si,"Sí"); label_out(*seg_no,"No")

    # Flechas verticales
    for i in range(len(rs)-1):
        p1=((rs[i][0]+rs[i][2])//2, rs[i][3]+SAFE)
        p2=((rs[i+1][0]+rs[i+1][2])//2, rs[i+1][1]-SAFE)
        arrow_down(d,p1,p2)
    for i in range(len(rn)-1):
        arrow_down(d, ((rn[i][0]+rn[i][2])//2, rn[i][3]+SAFE),
                      ((rn[i+1][0]+rn[i+1][2])//2, rn[i+1][1]-SAFE))

    # Retroalimentación externa
    rail_x=min(W-40, rs[-1][2]+int(retro_rail))
    start=(rs[-1][2]+SAFE, (rs[-1][1]+rs[-1][3])//2)
    mid1=(rail_x, start[1]); mid2=(rail_x, (r2[1]+r2[3])//2)
    end=(r2[2]+HEAD_CLEAR, (r2[1]+r2[3])//2)
    poly_arrow(d,[start,mid1,mid2,end])
    d.text((min(W-50,rail_x-10),(start[1]+mid2[1])//2),"Retroalimentación", font=FONT_SMALL, fill=BLUE, anchor="rm")

    # ===== Subproceso A: Nodos (SÍ-1 → … → SÍ-3) =====
    topA = r_fin[3] + 110
    d.text((W//2, topA-40), "Focalización por Nodos Demandantes (Proc. 1.4 – parte 1)",
           font=FONT_SUB, fill=BLUE, anchor="mm")
    bwA, gapA = int(anx_wA), int(anx_gapA); bhA = 92
    leftA, rightA = 80, W-80

    totA = 4*bwA + 3*gapA
    sxA = leftA + max(0,(rightA-leftA-totA)//2)
    yA_sup = topA + 120
    a_sup_txts=[a_sup1,a_sup2,a_sup3,a_sup4]; a_sup=[]
    for i,txt in enumerate(a_sup_txts):
        cxA = sxA + i*(bwA+gapA) + bwA//2
        r=[cxA-bwA//2, yA_sup-bhA//2, cxA+bwA//2, yA_sup+bhA//2]
        rrect(d,r); draw_centered(d,txt,r); a_sup.append(r)
    for i in range(3):
        arrow(d,(a_sup[i][2]+SAFE,yA_sup),(a_sup[i+1][0]-SAFE,yA_sup))

    y_railA = yA_sup - bhA//2 - 40
    poly_arrow(d,[(a_sup[0][0], y_railA),(a_sup[3][2], y_railA),(a_sup[3][2], yA_sup-bhA//2-6)], color=BLUE)

    a_inf_txts=[a_inf1,a_inf2,a_inf3,a_inf4,a_inf5,a_inf6]
    totA2 = 6*bwA + 5*gapA
    sxA2 = leftA + max(0,(rightA-leftA-totA2)//2)
    yA_inf = yA_sup + 220
    a_inf=[]
    for i,txt in enumerate(a_inf_txts):
        cxA = sxA2 + i*(bwA+gapA) + bwA//2
        r=[cxA-bwA//2, yA_inf-bhA//2, cxA+bwA//2, yA_inf+bhA//2]
        rrect(d,r); draw_centered(d,txt,r); a_inf.append(r)
    for i in range(5):
        arrow(d,(a_inf[i][2]+SAFE,yA_inf),(a_inf[i+1][0]-SAFE,yA_inf))
    arrow_down(d, ((a_sup[1][0]+a_sup[1][2])//2, a_sup[1][3]+SAFE), ((a_inf[0][0]+a_inf[0][2])//2, a_inf[0][1]-SAFE))
    arrow_down(d, ((a_sup[2][0]+a_sup[2][2])//2, a_sup[2][3]+SAFE), ((a_inf[2][0]+a_inf[2][2])//2, a_inf[2][1]-SAFE))

    # Conexiones de integración (SÍ-1 → Nodos → SÍ-3)
    start_si1=(rs[0][2]+SAFE, (rs[0][1]+rs[0][3])//2)
    targ_a1=(a_sup[0][0]-SAFE, yA_sup)
    poly_arrow(d,[start_si1,(start_si1[0]+60,start_si1[1]),(start_si1[0]+60,yA_sup),targ_a1], color=BLUE)
    end_a_last=(a_inf[-1][2]+SAFE, yA_inf)
    targ_s3=(rs[2][0]-SAFE, (rs[2][1]+rs[2][3])//2)
    poly_arrow(d,[end_a_last,(end_a_last[0]+60,yA_inf),(end_a_last[0]+60,targ_s3[1]),targ_s3], color=BLUE)

    # ===== Subproceso B: Reiterada (SÍ-6 → … → SÍ-7) =====
    topB = yA_inf + bhA//2 + 120
    d.text((W//2, topB-40), "Conducta delictiva reiterada", font=FONT_SUB, fill=BLUE, anchor="mm")
    lane_h, lane_gap, lane_w = 220, 24, 42
    lane_left, lane_right = 60, W-40
    for i,title in enumerate([lane1,lane2,lane3]):
        y0 = topB + i*(lane_h + lane_gap)
        y1 = y0 + lane_h
        d.rectangle([lane_left, y0, lane_right, y1], outline=BORDER, width=2)
        paste_vertical_label(img, (lane_left, y0, lane_left+lane_w, y1), title)

    bwB, gapB, bhB = int(anx_wB), int(anx_gapB), 92
    leftB = lane_left + lane_w + 24; rightB = lane_right - 24
    totalB = 4*bwB + 3*gapB
    sxB = leftB + max(0,(rightB-leftB-totalB)//2)
    yB_sup = topB + lane_h//2
    supB_txt=[r1,r2,r3,r4]; supB=[]
    for i,txt in enumerate(supB_txt):
        cx = sxB + i*(bwB+gapB) + bwB//2
        r=[cx-bwB//2, yB_sup-bhB//2, cx+bwB//2, yB_sup+bhB//2]
        rrect(d,r); draw_centered(d,txt,r); supB.append(r)
    for i in range(3):
        arrow(d,(supB[i][2]+SAFE,yB_sup),(supB[i+1][0]-SAFE,yB_sup))
    fin_w, fin_h = 120, 56
    cx_fin = supB[-1][2] + 90
    fin_box = [cx_fin-fin_w//2, yB_sup-fin_h//2, cx_fin+fin_w//2, yB_sup+fin_h//2]
    oval(d, fin_box); draw_centered(d,"FIN", fin_box)
    arrow(d, ((supB[-1][2]+SAFE), yB_sup), (fin_box[0]-SAFE, yB_sup))

    yB_mid = topB + lane_h + lane_gap + lane_h//2
    mid_x1 = sxB + bwB//2 + (bwB+gapB)*0.5
    mid_x2 = mid_x1 + bwB + gapB*1.2
    mid1=[int(mid_x1-bwB//2), yB_mid-bhB//2, int(mid_x1+bwB//2), yB_mid+bhB//2]
    mid2=[int(mid_x2-bwB//2), yB_mid-bhB//2, int(mid_x2+bwB//2), yB_mid+bhB//2]
    rrect(d,mid1); draw_centered(d,rm1,mid1)
    rrect(d,mid2); draw_centered(d,rm2,mid2)
    arrow(d,(mid1[2]+SAFE,yB_mid),(mid2[0]-SAFE,yB_mid))
    arrow_down(d, ((supB[2][0]+supB[2][2])//2, supB[2][3]+SAFE), ((mid1[0]+mid1[2])//2, mid1[1]-SAFE))
    arrow(d, ((mid2[0]+mid2[2])//2, mid2[1]-SAFE-20), ((supB[3][0]+supB[3][2])//2, supB[3][1]-SAFE-20))
    arrow(d, ((supB[3][0]+supB[3][2])//2, supB[3][1]-SAFE-20), ((supB[3][0]+supB[3][2])//2, supB[3][1]-SAFE))

    yB_inf = topB + 2*(lane_h + lane_gap) + lane_h//2
    inf_x = sxB + (bwB+gapB)*1.2
    inf=[int(inf_x-bwB//2), yB_inf-bhB//2, int(inf_x+bwB//2), yB_inf+bhB//2]
    rrect(d,inf); draw_centered(d,rb1,inf)
    arrow_down(d, ((mid1[0]+mid1[2])//2, mid1[3]+SAFE), ((inf[0]+inf[2])//2, inf[1]-SAFE))
    arrow(d, ((inf[0]+inf[2])//2, inf[1]-SAFE-20), ((mid2[0]+mid2[2])//2, mid2[3]+SAFE+20))
    arrow_down(d, ((mid2[0]+mid2[2])//2, mid2[3]+SAFE+20), ((mid2[0]+mid2[2])//2, mid2[3]+SAFE+22))

    # Conexiones de integración (SÍ-6 → Reiterada → SÍ-7)
    start_s6=(rs[5][2]+SAFE, (rs[5][1]+rs[5][3])//2)
    targ_b1=(supB[0][0]-SAFE, yB_sup)
    poly_arrow(d,[start_s6,(start_s6[0]+60,start_s6[1]),(start_s6[0]+60,yB_sup),targ_b1], color=BLUE)
    from_fin = (fin_box[2]+SAFE, yB_sup)
    targ_s7  = (rs[6][0]-SAFE, (rs[6][1]+rs[6][3])//2)
    poly_arrow(d,[from_fin,(from_fin[0]+60,yB_sup),(from_fin[0]+60,targ_s7[1]),targ_s7], color=BLUE)

    # --- Export PNG ---
    out=io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()

# ---------- Exportadores ----------
def make_pdf_from_png(png_bytes: bytes) -> bytes:
    img=PILImage.open(io.BytesIO(png_bytes)).convert("RGB")
    out=io.BytesIO(); img.save(out, format="PDF"); return out.getvalue()

def make_pptx(png_bytes: bytes) -> bytes:
    prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.2), Inches(0.2), width=Inches(9.6))
    out=io.BytesIO(); prs.save(out); return out.getvalue()

# ---------- Render & Descargas ----------
png_bytes=render_png()
pdf_bytes=make_pdf_from_png(png_bytes)
pptx_bytes=make_pptx(png_bytes)

st.subheader("Vista previa (proyecto integrado)")
st.image(png_bytes, use_column_width=True)

c1,c2,c3=st.columns(3)
with c1: st.download_button("⬇️ PNG", png_bytes, "mpgp_proyecto_integrado.png", "image/png")
with c2: st.download_button("⬇️ PDF", pdf_bytes, "mpgp_proyecto_integrado.pdf", "application/pdf")
with c3: st.download_button("⬇️ PPTX", pptx_bytes, "mpgp_proyecto_integrado.pptx",
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")

